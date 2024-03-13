import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def find_title_slide(prs, title):
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() == title:
                return i, shape
    return None, None

def insert_data():
    eng_number = eng_number_entry.get()
    engine_module = engine_module_entry.get()
    detail_pn = detail_pn_entry.get()
    assy_pn = assy_pn_entry.get()
    vendor = vendor_entry.get()
    qn = qn_entry.get()
    sn = sn_entry.get()
    line_item = line_item_entry.get()
    
    # Concatenate all input data into one string
    user_input = f"ENG Number: {eng_number}\n" \
                 f"Engine Module: {engine_module}\n" \
                 f"Detail P/N: {detail_pn}\n" \
                 f"Assy P/N: {assy_pn}\n" \
                 f"Vendor: {vendor}\n" \
                 f"QN: {qn}\n" \
                 f"S/N: {sn}\n" \
                 f"Line Item: {line_item}"
    
    selected_file = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    
    if selected_file:
        prs = Presentation(selected_file)
        
        # Insert user input into the first slide
        slide = prs.slides[0]
        left = Inches(1)
        top = Inches(4)
        width = Inches(8)
        height = Inches(4)
        
        textbox_user_input = slide.shapes.add_textbox(left, top, width, height)
        text_frame_user_input = textbox_user_input.text_frame
        
        # Set text alignment to left
        text_frame_user_input.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add user input as a single paragraph
        p_user_input = text_frame_user_input.add_paragraph()
        p_user_input.text = user_input
        p_user_input.space_after = Pt(10)  # Adjust spacing between paragraphs
        
        # Update "Interim 502 Recommendation" slide
        title_slide_index, title_shape = find_title_slide(prs, "Interim 502 Recommendation")
        if title_slide_index is not None:
            # Calculate position for QN number on top of the title shape
            qn_x = title_shape.left
            qn_y = title_shape.top - Inches(0.5)  # Adjust this value to set the distance from the title
            
            # Add QN number
            textbox_qn = prs.slides[title_slide_index].shapes.add_textbox(qn_x, qn_y, Inches(2), Inches(0.5))
            text_frame_qn = textbox_qn.text_frame
            p_qn = text_frame_qn.add_paragraph()
            p_qn.text = f"QN: {qn}"
            p_qn.font.size = Pt(18)
        else:
            messagebox.showwarning("Warning", "Title 'Interim 502 Recommendation' not found in the presentation.")
        
        # Update "Design PL3 Approval" slide
        title_slide_index, title_shape = find_title_slide(prs, "Design PL3 Approval")
        if title_slide_index is not None:
            # Calculate position for QN number on top of the title shape
            qn_x = title_shape.left
            qn_y = title_shape.top - Inches(0.5)  # Adjust this value to set the distance from the title
            
            # Add QN number
            textbox_qn = prs.slides[title_slide_index].shapes.add_textbox(qn_x, qn_y, Inches(2), Inches(0.5))
            text_frame_qn = textbox_qn.text_frame
            p_qn = text_frame_qn.add_paragraph()
            p_qn.text = f"QN: {qn}"
            p_qn.font.size = Pt(18)
        else:
            messagebox.showwarning("Warning", "Title 'Design PL3 Approval' not found in the presentation.")
        
        # Update "Engine Cross Section" slide
        title_slide_index, title_shape = find_title_slide(prs, "Engine Cross Section")
        if title_slide_index is not None:
            # Calculate position for ENG number on top of the title shape
            eng_number_x = title_shape.left
            eng_number_y = title_shape.top - Inches(0.5)  # Adjust this value to set the distance from the title
            
            # Add ENG number
            textbox_eng_number = prs.slides[title_slide_index].shapes.add_textbox(eng_number_x, eng_number_y, Inches(2), Inches(0.5))
            text_frame_eng_number = textbox_eng_number.text_frame
            p_eng_number = text_frame_eng_number.add_paragraph()
            p_eng_number.text = f"ENG Number: {eng_number}"
            p_eng_number.font.size = Pt(18)
        else:
            messagebox.showwarning("Warning", "Title 'Engine Cross Section' not found in the presentation.")
        
        # Update "Interfacing Parts (3D View)" slide
        title_slide_index, title_shape = find_title_slide(prs, "Interfacing Parts (3D View)")
        if title_slide_index is not None:
            # Calculate position for ENG number on top of the title shape
            eng_number_x = title_shape.left
            eng_number_y = title_shape.top - Inches(0.5)  # Adjust this value to set the distance from the title
            
            # Add ENG number
            textbox_eng_number = prs.slides[title_slide_index].shapes.add_textbox(eng_number_x, eng_number_y, Inches(2), Inches(0.5))
            text_frame_eng_number = textbox_eng_number.text_frame
            p_eng_number = text_frame_eng_number.add_paragraph()
            p_eng_number.text = f"ENG Number: {eng_number}"
            p_eng_number.font.size = Pt(18)
        else:
            messagebox.showwarning("Warning", "Title 'Interfacing Parts (3D View)' not found in the presentation.")
        
        prs.save(selected_file)
        messagebox.showinfo("Success", "Data inserted into PowerPoint successfully!")

# GUI setup
root = tk.Tk()
root.title("Insert Data into PowerPoint")

tk.Label(root, text="ENG Number:").grid(row=0, column=0)
eng_number_entry = tk.Entry(root)
eng_number_entry.grid(row=0, column=1)

tk.Label(root, text="Engine Module:").grid(row=1, column=0)
engine_module_entry = tk.Entry(root)
engine_module_entry.grid(row=1, column=1)

tk.Label(root, text="Detail P/N:").grid(row=2, column=0)
detail_pn_entry = tk.Entry(root)
detail_pn_entry.grid(row=2, column=1)

tk.Label(root, text="Assy P/N:").grid(row=3, column=0)
assy_pn_entry = tk.Entry(root)
assy_pn_entry.grid(row=3, column=1)

tk.Label(root, text="Vendor:").grid(row=4, column=0)
vendor_entry = tk.Entry(root)
vendor_entry.grid(row=4, column=1)

tk.Label(root, text="QN:").grid(row=5, column=0)
qn_entry = tk.Entry(root)
qn_entry.grid(row=5, column=1)

tk.Label(root, text="S/N:").grid(row=6, column=0)
sn_entry = tk.Entry(root)
sn_entry.grid(row=6, column=1)

tk.Label(root, text="Line Item:").grid(row=7, column=0)  # New label for Line Item
line_item_entry = tk.Entry(root)  # New entry field for Line Item
line_item_entry.grid(row=7, column=1)  # Adjust the grid layout

insert_button = tk.Button(root, text="Insert Data into PowerPoint", command=insert_data)
insert_button.grid(row=8, columnspan=2)  # Adjust the grid layout

root.mainloop()
