import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def insert_data():
    eng_number = eng_number_entry.get()
    engine_module = engine_module_entry.get()
    detail_pn = detail_pn_entry.get()
    assy_pn = assy_pn_entry.get()
    vendor = vendor_entry.get()
    qn = qn_entry.get()
    sn = sn_entry.get()
    line_item = line_item_entry.get()  # Get Line Item value
    
    data = {'ENG Number': eng_number,
            'Engine Module': engine_module,
            'Detail P/N': detail_pn,
            'Assy P/N': assy_pn,
            'Vendor': vendor,
            'QN': qn,
            'S/N': sn,
            'Line Item': line_item}  # Include Line Item in the data dictionary
    
    selected_file = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    
    if selected_file:
        prs = Presentation(selected_file)
        slide = prs.slides[0]  # Assuming the data should be inserted into the first slide
        
        # Calculate the position for the textbox
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Set text alignment to left
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        for key, value in data.items():
            # Add each key-value pair as a separate paragraph
            p = text_frame.add_paragraph()
            p.text = f"{key}: {value}"
            p.space_after = Pt(10)  # Adjust spacing between paragraphs
        
        prs.save(selected_file)
        messagebox.showinfo("Success", "Data inserted into PowerPoint successfully!")

# GUI setup
root = tk.Tk()
root.title("Insert Data into PowerPoint")

tk.Label(root, text="ENG Number:").grid(row=0, column=0)
eng_number_entry = tk.Entry(root)
eng_number_entry.grid(row=0, column=1)

tk.Label(root, text="Engine Module:").grid(row=1, column=0)
engine_module_entry = tk.Entry(root)
engine_module_entry.grid(row=1, column=1)

tk.Label(root, text="Detail P/N:").grid(row=2, column=0)
detail_pn_entry = tk.Entry(root)
detail_pn_entry.grid(row=2, column=1)

tk.Label(root, text="Assy P/N:").grid(row=3, column=0)
assy_pn_entry = tk.Entry(root)
assy_pn_entry.grid(row=3, column=1)

tk.Label(root, text="Vendor:").grid(row=4, column=0)
vendor_entry = tk.Entry(root)
vendor_entry.grid(row=4, column=1)

tk.Label(root, text="QN:").grid(row=5, column=0)
qn_entry = tk.Entry(root)
qn_entry.grid(row=5, column=1)

tk.Label(root, text="S/N:").grid(row=6, column=0)
sn_entry = tk.Entry(root)
sn_entry.grid(row=6, column=1)

tk.Label(root, text="Line Item:").grid(row=7, column=0)  # New label for Line Item
line_item_entry = tk.Entry(root)  # New entry field for Line Item
line_item_entry.grid(row=7, column=1)  # Adjust the grid layout

insert_button = tk.Button(root, text="Insert Data into PowerPoint", command=insert_data)
insert_button.grid(row=8, columnspan=2)  # Adjust the grid layout

root.mainloop()
