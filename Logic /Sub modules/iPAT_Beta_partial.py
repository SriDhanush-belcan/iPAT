   
from tkinter import *
from tkinter.ttk import Label
from tkinter import Tk
from tkinter.simpledialog import askinteger
from tkinter import messagebox
from functools import partial
import tkinter as tk
from tkinter import filedialog
from PIL import Image, ImageTk
import os
from pptx import Presentation  
from pptx.util import Inches
from pptx.util import Pt 
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.util import Inches, Pt
import codecs
from pptx.enum.text import PP_ALIGN

master = Tk()
# Adjust size
master.geometry("590x550")
#Srikanth code
w = Canvas(master, width=365, height=550)
#master = Canvas(master, width=400, height=300)
w.place(x=0,y=0)
w.configure(bg='#8eb9d8')  # , borderwidth
master['background']='#1d5075'
master.title('iPAT')

def printDetails(usernameEntry) :
    usernameText = usernameEntry.get()
    print("user entered :", usernameText)
    return

def save_to_file(data):
    with open("user_input.txt", "w") as file:
    #with open("user_input.pptx", "w") as file:
        file.write(data)


def add_page_numbers(prs):
    for slide in prs.slides:
        # Add a footer to the slide
        footer_shape = slide.shapes.add_textbox(left=Inches(8), top=Inches(7.2), width=Inches(2), height=Inches(0.5))
        text_frame = footer_shape.text_frame
        text_frame.clear()  # Clear any existing text
        p = text_frame.paragraphs[0]
        
        # Set the text to the slide number
        p.text = f"{prs.slides.index(slide) + 1}"
        
        # Set font properties
        run = p.runs[0]
        font = run.font
        font.size = Pt(8)
        font.bold = False
        font.italic = False
        font.color.rgb = RGBColor(0, 0, 0)  # Black color
        font.name = 'Arial'
        text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Add the additional text
        additional_text = "Pratt & Whitney PROPRIETARY Data. Releasable to P&W USA, Belcan USA, and Belcan India Only. Subject to the export control restrictions on the title page of this document"
        additional_textbox = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(7.1), width=Inches(9), height=Inches(0.5))
        additional_text_frame = additional_textbox.text_frame
        additional_text_frame.clear()
        additional_p = additional_text_frame.paragraphs[0]
        additional_p.text = additional_text
        additional_run = additional_p.runs[0]
        additional_font = additional_run.font
        additional_font.size = Pt(8)
        additional_font.bold = False
        additional_font.italic = False
        additional_font.color.rgb = RGBColor(0, 0, 0)  # Black color
        additional_font.name = 'Arial'
        additional_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Load the existing presentation
#prs = Presentation('Output.pptx')

# Call the function to add page numbers to all slides
#add_page_numbers(prs)

def Ref_ok():
    QN1Value = QN1Val.get()
    PartNoVal = PartNo.get()
    AssyNoVal = AssyNo.get()
    LIVal = LI1Val.get()        
    PartNameVal = PartName.get()
    SerialNoVal = SerialNo.get()
    VendorNVal = VendorN.get()

    data = f"QN#          : {QN1Value}\n" \
           f"Part No      : {PartNoVal}\n" \
           f"Assy No      : {AssyNoVal}\n" \
           f"LI           : {LIVal}\n" \
           f"Part Name    : {PartNameVal}\n" \
           f"Serial No    : {SerialNoVal}\n" \
           f"Vendor       : {VendorNVal}\n"

    save_to_file(data)
    
    messagebox.showinfo("File Saved", "Input data has been saved to user_input.txt")

    # Insert data into the table on QN RESOLUTION DATA slide
    #insert_data_into_ppt(QN1Value, PartNoVal, AssyNoVal, LIVal, PartNameVal, SerialNoVal, VendorNVal)
    insert_data_into_ppt(PartNoVal, AssyNoVal, VendorNVal, QN1Value, SerialNoVal)


def qexit():
    master.destroy()

#import aspose.slides as slides
from pptx import Presentation
from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN
# importing Inches from the util class of the pptx library  
from pptx.util import Inches  
imgPath = "logo.png"  

#def insert_data_into_ppt(QN1Value, PartNoVal, AssyNoVal, LIVal, PartNameVal, SerialNoVal, VendorNVal):
def insert_data_into_ppt(PartNoVal, AssyNoVal, VendorNVal, QN1Value, SerialNoVal):
    # Creating a Presentation
    myPPT = Presentation()

    # select a layout for the presentation
    slideLayout = myPPT.slide_layouts[5]
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'QN RESOLUTION DATA'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)
    

    # left, top, width, height
    shape1 = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.7), Inches(4), Inches(2.7))
    #shape1 = rectangle
    shape1.shadow.inherit = False
    fill=shape1.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(194, 233, 203)

    # Inserting text boxes to display data
    data_text = f"ENG Number, Engine Module\n"\
                f"Details P/N : {PartNoVal}\n"\
                f"Assy No : {AssyNoVal}\n"\
                f"Vendor : {VendorNVal}\n"\
                f"QN : {QN1Value}\n"\
                f"Serial No : {SerialNoVal}\n"\
                f" RESOLUTION DATA"
    

    

    left = Inches(5.2)
    top = Inches(1.7)
    width = Inches(4)
    height = Inches(2.6)
    text_box = shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    
    # Add each piece of data as a separate paragraph with proper indentation
    #for line in data_text.split('\n'):
    for line in data_text.split('\n'):
       
        p = text_frame.add_paragraph()
        p.text = line 
        p.level = 0  # Adjust the indentation level as needed for proper alignment
        title_shape = p
        #title_shape.text_frame.font.name = 'Arial'
        title_shape.font.name = 'Arial'
        #title_shape.text.font.name = 'Arial'
        title_shape.font.size = Pt(20)
        title_shape.font.bold = True
        title_shape.alignment = PP_ALIGN.CENTER      

    left = Inches(5.1)
    top = Inches(4.6)
    width = Inches(4.5)
    height = Inches(1.2)
    text_box = shapes.add_textbox(left, top, width, height)
    tf1 = text_box.text_frame
    p = tf1.paragraphs[0]
    #p = tf1.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()

    #tf1.text = "Module Part Family, Hot/Cold Section Engineering \n XXXXX- Belcan Design \n XXXXX- Belcan Design \n XXXXX- Belcan Design \n XXXXXX P&W Design PL3"
    run.text = "Module Part Family, Hot/Cold Section Engineering \n XXXXX- Belcan Design \n XXXXX- Belcan Design \n XXXXX- Belcan Design \n XXXXXX P&W Design PL3"
    #font = tf1.font
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    #font.bold = True
    font.italic = None  # cause value to be inherited from theme


    # left, top, width, height
    shape2 = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.7), Inches(4), Inches(4))
    #shape1 = rectangle
    shape2.shadow.inherit = False
    fill=shape2.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(194, 233, 203)

    left = Inches(0.6)
    top = Inches(1.7)
    width = Inches(4)
    height = Inches(4)
    text_box = shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame

    left = Inches(0.5)
    top = Inches(5.8)
    width = Inches(9)
    height = Inches(1.2)
    text_box = shapes.add_textbox(left, top, width, height)
    tf1 = text_box.text_frame
    tf1 = tf1.paragraphs[0]
    tf1.text = "Title_Slide_Marking\n"
    font = tf1.font
    font.name = 'Arial'
    font.size = Pt(16)
    font.bold = True
    font.italic = True  # cause value to be inherited from theme

    #shape1.set_footer_text("New footer text")
    # save file 
    myPPT.save('Output.pptx')

    # Load the existing presentation
    prs = Presentation('Output.pptx')

    # Call the function to add page numbers to all slides
    add_page_numbers(prs)

QN1Val = StringVar()
PartNo = StringVar()
AssyNo = StringVar()
LI1Val = StringVar()
PartName = StringVar()
SerialNo = StringVar()
VendorN = StringVar()

# Text Button1
txtbxqn = Label(master, text='QN#                            :').place(x=8,y=10)
# Entry for user input
txtbxqn1 = Entry(master, textvariable=QN1Val).place(x=150,y=10)

# Text Button2
txtbxPartno = Label(master, text='Part No                       :').place(x=8,y=40)
# Entry for user input
txtbxPartno1 = Entry(master, textvariable=PartNo).place(x=150,y=40)

# Text Button3
txtbxassyno = Label(master, text='Assy No                      :').place(x=8,y=70)
txtbxassyno1= Entry(master, textvariable=AssyNo).place(x=150,y=70)

# Text Button4
txtbxmqi = Label(master, text='LI Value                        :').place(x=8,y=100)
txtbxmqi1 = Entry(textvariable=LI1Val).place(x=150,y=100)

# Text Button5
txtbxpartname = Label(master, text='Part Name                  :').place(x=8,y=130)
txtbxpartname1 = Entry(master, textvariable=PartName).place(x=150,y=130)

# Text Button6
txtbxserialn = Label(master, text='Serial No                     :').place(x=8,y=160)
txtbxserialn1 = Entry(master, textvariable=SerialNo).place(x=150,y=160)

# iPAT Display
usernameLabel= Label(master, text='i\nP\nA\nT',font=('Arial',60),foreground='#5783a3',background='#1d5075').place(x=520,y=30)

# Text Button7
txtbxvendor = Label(master, text='Vendor                        :').place(x=8,y=190)
txtbxvendor1 = Entry(master, textvariable=VendorN).place(x=150,y=190)

#browse button1
pushbuttnsap = Label(master, text='SAP Package PPT      :').place(x=8,y=250)
# Function to update the Listbox with folder contents
def browse_folder():
    folder_path = filedialog.askdirectory()  # Open a folder selection dialog
    if folder_path:
        folder_contents.delete(0, tk.END)  # Clear the Listbox
        for item in os.listdir(folder_path):
            folder_contents.insert(tk.END, item)  # Insert folder contents into Listbox

# Create a button to browse for a folder
pushbuttnsap1 = tk.Button(master, text="Browse SAP Package", command=browse_folder, activebackground="#1d5075",activeforeground="#8eb9d8",bd=3).place(x=150,y=246)

#browse button2
pushbuttnvendor = Label(master, text='Vendor info PPT        :').place(x=8,y=300)
# Function to update the Listbox with folder contents
def browse_folder():
    folder_path = filedialog.askdirectory()  # Open a folder selection dialog
    if folder_path:
        folder_contents.delete(0, tk.END)  # Clear the Listbox
        for item in os.listdir(folder_path):
            folder_contents.insert(tk.END, item)  # Insert folder contents into Listbox

# Create a button to browse for a folder
pushbuttnvendor1 = tk.Button(master, text=" Browse Vendor Info ", command=browse_folder, activebackground="#1d5075",activeforeground="#8eb9d8",bd=3).place(x=150,y=296)

#Select List button
selectlist = Label(master, text='Engine Model            :').place(x=8,y=350)
# Create the list of options 
options_list = ["PW1100G", "PW1500G", "PW1900G"] 
#print(options_list)
# Variable to keep track of the option 
# selected in OptionMenu 
value_inside = tk.StringVar(master) 
  
# Set the default value of the variable 
value_inside.set("Select Engine Model") 
  
# Create the optionmenu widget and passing  
# the options_list and value_inside to it. 
question_menu = tk.OptionMenu(master, value_inside, *options_list).place(x=150,y=346)

# Function to print the submitted option-- testing purpose 
def print_answers(): 
    print("Selected Option: {}".format(value_inside.get())) 
    return None


# Function to update the Listbox with folder contents
def browse_folder():
    folder_path = filedialog.askdirectory()  # Open a folder selection dialog
    if folder_path:
        folder_contents.delete(0, tk.END)  # Clear the Listbox
        for item in os.listdir(folder_path):
            folder_contents.insert(tk.END, item)  # Insert folder contents into Listbox


# Create a button to browse for a folder
btnTotal=Button(fg="black",font=('ariel' ,10,'bold'),height=2, width=8, text="OK", command=Ref_ok)
btnTotal.place(x=420,y=496)

btnexit=Button(fg="black",font=('ariel' ,10,'bold'),height=2, width=8, text="Generate", command=qexit)
btnexit.place(x=496,y=496)

# Execute tkinter
master.mainloop()

##################################
####### Summary Table ############
##################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[5]  
mySlide = myPPT.slides.add_slide(slideLayout)  
shapes = mySlide.shapes  

#font = shapes.font
shapes.title.text = 'Summary Table'
#shapes.title.text = 'Prior History by MQI or Damage Code'
title_shape = shapes.title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

# specifying the rows, columns, and other measurements  
LI1Val1 = LI1Val.get()
print(LI1Val1)
rows = int(LI1Val1)
rows1 = rows + 1
#rows = 3
cols = 5  
left = Inches(0.2)
top = Inches(1.6)  
width = Inches(9.6)  
height = Inches(0.8)  
  
# using the add_table() method  
myTable = shapes.add_table(rows1, cols, left, top, width, height).table 

# writing column headings  
myTable.cell(0, 0).text = 'Line Item'  
myTable.cell(0, 1).text = 'MQI'  
myTable.cell(0, 2).text = 'Location'  
myTable.cell(0, 3).text = 'Description'
myTable.cell(0, 4).text = 'Recommendation'  

# writing body cells  
myTable.cell(1, 0).text = '1'  
myTable.cell(1, 1).text = 'Ram'
myTable.cell(1, 2).text = 'Enter manually'
myTable.cell(1, 3).text = 'Enter manually'
myTable.cell(1, 4).text = 'Enter manually'

# Set Font size and type in table 
def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


for cell in iter_cells(myTable):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)
            run.font.name = 'Arial'

# Add slide number
add_page_numbers(myPPT)

# saving the PPT file  
myPPT.save('Output.pptx')  

#############################################
########## Defect Description ###############
#############################################
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Defect Description'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape1 = rectangle
    shape1.shadow.inherit = False
    fill=shape1.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
  
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.3)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches() 
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.52), Inches(10), Inches(0.30)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Enter manually"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)
    
    # Add slide number
    add_page_numbers(myPPT)
    # saving the PPT file  
    myPPT.save('Output.pptx')  

    ##############################################
    ########### Defect Location ##################
    ##############################################
    # creating the slide  
    myPPT = Presentation('Output.pptx')  
    #for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Defect Location'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.3)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches() 
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.52), Inches(10), Inches(0.30)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": B/P Location of defect"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')  

##########################################
#### Defect MQI Location #############
##########################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Defect MQI Location'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.3)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches() 
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.52), Inches(10), Inches(0.30)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": B/P Location of defect"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################
#### Disposition & Approvals #############
##########################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Disposition & Approvals'
title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
subtitle.text = '\n'

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

# Add slide number
add_page_numbers(myPPT)

# saving the PPT file  
myPPT.save('Output.pptx')  

#####################################################
######### Interim 502 Recommendation ################
#####################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes
    qn_value = QN1Val.get()
    shapes.title.text = 'Interim 502 Recommendation XXXX'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(4), Inches(1.0), Inches(2.5), Inches(0.35)) 
    # creating textFrames 
    tf = txBox1.text_frame 
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    run.text = f"QN {qn_value}"  # Update the QN value
    font = run.font
    font.name = 'Arial'
    font.size = Pt(28)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches() 
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.52), Inches(10), Inches(0.30)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": XXXXX"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)
       
    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0),Inches(9),Inches(0.3))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(217, 217, 217)
    #fill.fore_color.rgb=RGBColor(188, 190, 200)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox3 = mySlide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.3)) 
    # creating textFrames 
    text_frame = txBox3.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ", MQI XXXX:"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(16, 17, 17)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')  

#################################################
#########   Design PL3 Approval   ###############
#################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes
    qn_value = QN1Val.get()  
    shapes.title.text = 'Design PL3 Approval'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(4), Inches(0.2), Inches(2.5), Inches(0.35)) 
    # creating textFrames 
    tf = txBox1.text_frame 
    #tf.text = "QN XXXXXXXXX "
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    #run.text = "LINE ITEM" + "-" + str(i+1)
    run.text = f"QN {qn_value}"  # Update the QN value
    font = run.font
    font.name = 'Arial'
    font.size = Pt(28)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
 
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1) + ":"
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.35)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ":"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################
######## Design Assessment  ##############
##########################################
# creating the slide  
myPPT = Presentation('Output.pptx')
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Design Assessment'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
subtitle.text = '\n'

# Add slide number
add_page_numbers(myPPT)

# saving the PPT file  
myPPT.save('Output.pptx')  

#########################################################
########### Design Analysis Summary ####################
#########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')
for i in range(rows): 
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Design Analysis Summary'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(4), Inches(0.4)) 

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.3)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": XXXXX"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

########################################################
################  Engine Cross Section  ################
########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Engine Cross Section'
    # Get the engine module value from the GUI input
    engine_module = value_inside.get()

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.4)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Defect location in Engine Cross Section"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(3.5), Inches(2), Inches(3.5), Inches(0.5)) 
    # creating textFrames 
    tf = txBox1.text_frame 
    #tf.text = "QN XXXXXXXXX "
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    # Set the text frame with the engine module value
    run.text = f"Engine Model: {engine_module}"

    font = run.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################################
#############  Interfacing Parts  ########################
##########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Interfacing Parts'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        
    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.35)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Enter Manually"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

    ###############################################################
    ################ Interfacing Parts (3D View) ##################
    ###############################################################
    # creating the slide  
    myPPT = Presentation('Output.pptx')  
    #for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Interfacing Parts (3D View)'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        
    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.35)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Enter Manually"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Get the engine module value from the GUI input
    engine_module = value_inside.get()

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(6), Inches(2.1), Inches(3.5), Inches(0.4)) 
    # creating textFrames 
    tf = txBox1.text_frame 
    #tf.text = "QN XXXXXXXXX "
    p1 = tf.paragraphs[0]
    run = p1.add_run()
    # Set the text frame with the engine module value
    run.text = f"Engine Model: {engine_module}"

    font = run.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

######################################################
##########  SPEC: XXXX XXXXXX ########################
######################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
#for i in range(rows):
slideLayout = myPPT.slide_layouts[5]  
mySlide = myPPT.slides.add_slide(slideLayout)  
shapes = mySlide.shapes  
    
#font = shapes.font
shapes.title.text = 'SPEC: XXX XXXXXX'
#slide.shapes.title.font.name = 'Arial'
#shapes.title.name = 'Arial'
title_shape = shapes.title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
#rectangle.color.rgb = RGBColor(55, 238, 95)
shape = rectangle
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(14, 127, 38)

# creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.35)) 
text_frame = txBox.text_frame
text_frame.clear()  # not necessary for newly-created shape
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = "LINE ITEM" + "-" + str(i) + "," + str(i+1)
font = run.font
#font.name = 'Calibri'
font.name = 'Arial'
font.size = Pt(12)
font.bold = True
font.italic = None  # cause value to be inherited from theme
#font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

# Insert Body Line
line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.95), Inches(0), Inches(1.95))
line3.line.color.rgb = RGBColor(0,  112,  192)
line3.line.width = Pt(1.0)

line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.98), Inches(0), Inches(1.98))
line4.line.color.rgb = RGBColor(0,  112,  192)
line4.line.width = Pt(1.0)
    
# creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.5)) 
# creating textFrames 
text_frame = txBox1.text_frame
text_frame.clear()  # not necessary for newly-created shape
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = "Line Item" + "-" + str(i) + ": Enter Manually \nLine Item" + "-" + str(i+1) + ": Enter Manually"
font = run.font
font.name = 'Arial'
font.size = Pt(12)
font.bold = True
font.italic = None
font.color.rgb = RGBColor(0,  112,  192)

# Add slide number
add_page_numbers(myPPT)
 
# saving the PPT file  
myPPT.save('Output.pptx')

###########################################################
#################  Vendor Information  ####################
###########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Vendor Information'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
subtitle.text = '\n'

# Add slide number
add_page_numbers(myPPT)

# saving the PPT file  
myPPT.save('Output.pptx')  

#######################################################
########   Vendor Supplied Information ################
#######################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Vendor Supplied Information'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##############################################################
############## Engine Manual   ###############################
##############################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Engine Manual'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
subtitle.text = '\n'

# Add slide number
add_page_numbers(myPPT)

# saving the PPT file  
myPPT.save('Output.pptx')  

########################################################
############ Engine Manual #############################
########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Engine Manual'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.35)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Enter Manually"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

####################################################################
########### Root Cause, Corrective Action & Prior History ##########
####################################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Root Cause, Corrective Action & Prior History'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# Insert Line
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.shadow.inherit = False
line1.line.color.rgb = RGBColor(204, 0, 0)
line1.line.width = Pt(4.0)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
line2.shadow.inherit = False
line2.line.color.rgb = RGBColor(204, 0, 0)
line2.line.width = Pt(1.0)

# specifying the values of the parameters for the add_picture() method 
# left, top, width, height 
myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
#title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = '\n'

# Add slide number
add_page_numbers(myPPT)

# saving the PPT file  
myPPT.save('Output.pptx')

#########################################################
########### Root Cause & Corrective Action ##############
#########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Root Cause & Corrective Action'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    
    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.4)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Root cause & Corrective Action"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

#########################################################
########### Prior History by S/N  #######################
#########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Prior History by S/N'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert shape
    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    #font.name = 'Calibri'
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.35)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Enter Manually"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################################
#########  Prior History by MQI or Damage Code ###########
##########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Prior History by MQI or Damage Code'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert Line
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.shadow.inherit = False
    line1.line.color.rgb = RGBColor(204, 0, 0)
    line1.line.width = Pt(4.0)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.55), Inches(0), Inches(1.55))
    line2.shadow.inherit = False
    line2.line.color.rgb = RGBColor(204, 0, 0)
    line2.line.width = Pt(1.0)

    # specifying the values of the parameters for the add_picture() method 
    # left, top, width, height 
    myImage = mySlide.shapes.add_picture(imgPath, Inches(0), Inches(0.1), Inches(1.5), Inches(1.3)) 

    # Insert shape
    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.2),Inches(0.35))
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox = mySlide.shapes.add_textbox(Inches(8), Inches(1.1), Inches(1.2), Inches(0.4)) 
    text_frame = txBox.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "LINE ITEM" + "-" + str(i+1)
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    # Insert Body Line
    line3= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line3.line.color.rgb = RGBColor(0,  112,  192)
    line3.line.width = Pt(1.0)

    line4= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.88), Inches(0), Inches(1.88))
    line4.line.color.rgb = RGBColor(0,  112,  192)
    line4.line.width = Pt(1.0)
    
    # creating textBox, In terms of left = Inches(), top = Inches(), width =  Inches(),  height = Inches()
    txBox1 = mySlide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(0.35)) 
    # creating textFrames 
    text_frame = txBox1.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Line Item" + "-" + str(i+1) + ": Enter Manually"
    font = run.font
    font.name = 'Arial'
    font.size = Pt(12)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0,  112,  192)

    # Add slide number
    add_page_numbers(myPPT)

    # saving the PPT file  
    myPPT.save('Output.pptx')

