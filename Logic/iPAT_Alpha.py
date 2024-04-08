from tkinter import *
from tkinter.ttk import Label
from tkinter import Tk
from tkinter.simpledialog import askinteger
from tkinter import messagebox
from functools import partial
import tkinter as tk
from tkinter import filedialog
from PIL import Image, ImageTk
import os
from pptx import Presentation  
from pptx.util import Inches
from pptx.util import Pt 
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.util import Inches, Pt
import codecs
from pptx.enum.text import PP_ALIGN

master = Tk()
# Adjust size
master.geometry("590x550")
#Srikanth code
w = Canvas(master, width=365, height=550)
#master = Canvas(master, width=400, height=300)
w.place(x=0,y=0)
#w.pack(pady = 5)
#w.pack(side='left',expand = True, fill='both')

w.configure(bg='#8eb9d8')  # , borderwidth
# 2nd rect
#w1 = Canvas(master, width=50, height=50)
#master = Canvas(master, width=400, height=300)
#w1.place(x=4,y=225)
#w1.configure(bg='#8eb9d8')
# Set window color
#master.configure(bg='blue')

master['background']='#1d5075'

master.title('iPAT')

def printDetails(usernameEntry) :
    usernameText = usernameEntry.get()
    print("user entered :", usernameText)
    return


def save_to_file(data):
    with open("user_input.txt", "w") as file:
    #with open("user_input.pptx", "w") as file:
        file.write(data)

def Ref_ok():
    QN1Value = QN1Val.get()
    PartNoVal = PartNo.get()
    AssyNoVal = AssyNo.get()
    LIVal = LI1Val.get()        
    PartNameVal = PartName.get()
    SerialNoVal = SerialNo.get()
    VendorNVal = VendorN.get()

    data = f"QN#          : {QN1Value}\n" \
           f"Part No      : {PartNoVal}\n" \
           f"Assy No      : {AssyNoVal}\n" \
           f"LI           : {LIVal}\n" \
           f"Part Name    : {PartNameVal}\n" \
           f"Serial No    : {SerialNoVal}\n" \
           f"Vendor       : {VendorNVal}\n"

    save_to_file(data)

    messagebox.showinfo("File Saved", "Input data has been saved to user_input.txt")
    #messagebox.showinfo("File Saved", "Input data has been saved to user_input.pptx")

    # Insert data into the table on QN RESOLUTION DATA slide
    insert_data_into_ppt(QN1Value, PartNoVal, AssyNoVal, LIVal, PartNameVal, SerialNoVal, VendorNVal)

def qexit():
    master.destroy()

def insert_data_into_ppt(QN1Value, PartNoVal, AssyNoVal, LIVal, PartNameVal, SerialNoVal, VendorNVal):
    # Creating a Presentation
    myPPT = Presentation()

    # select a layout for the presentation
    slideLayout = myPPT.slide_layouts[5]
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'QN RESOLUTION DATA'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

        
# Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)
    
    # Inserting text boxes to display data
    data_text = f"QN#               : {QN1Value}\n"\
                f"Part No          : {PartNoVal}\n"\
                f"Assy No         : {AssyNoVal}\n"\
                f"LI                    :  {LIVal}\n"\
                f"Part Name     : {PartNameVal}\n"\
                f"Serial No        : {SerialNoVal}\n"\
                f"Vendor           : {VendorNVal}\n"
    
    left = Inches(5)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(6)
    text_box = shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    
    # Add each piece of data as a separate paragraph with proper indentation
    for line in data_text.split('\n'):
        p = text_frame.add_paragraph()
        p.text = line 
        p.level = 0  # Adjust the indentation level as needed for proper alignment

    # save file 
    myPPT.save('Output.pptx')

QN1Val = StringVar()
PartNo = StringVar()
AssyNo = StringVar()
LI1Val = StringVar()
PartName = StringVar()
SerialNo = StringVar()
VendorN = StringVar()

# Text Button1
txtbxqn = Label(master, text='QN#                            :').place(x=8,y=10)
# Entry for user input
txtbxqn1 = Entry(master, textvariable=QN1Val).place(x=150,y=10)

# Text Button2
txtbxPartno = Label(master, text='Part No                       :').place(x=8,y=40)
# Entry for user input
txtbxPartno1 = Entry(master, textvariable=PartNo).place(x=150,y=40)

# Text Button3
txtbxassyno = Label(master, text='Assy No                      :').place(x=8,y=70)
txtbxassyno1= Entry(master, textvariable=AssyNo).place(x=150,y=70)

# Text Button4
txtbxmqi = Label(master, text='LI Value                        :').place(x=8,y=100)
#txtbxmqi1 = Entry(master).place(x=150,y=100)
txtbxmqi1 = Entry(textvariable=LI1Val).place(x=150,y=100)


# Text Button5
txtbxpartname = Label(master, text='Part Name                  :').place(x=8,y=130)
txtbxpartname1 = Entry(master, textvariable=PartName).place(x=150,y=130)

# Text Button6
txtbxserialn = Label(master, text='Serial No                     :').place(x=8,y=160)
txtbxserialn1 = Entry(master, textvariable=SerialNo).place(x=150,y=160)

# iPAT Display
usernameLabel= Label(master, text='i\nP\nA\nT',font=('Arial',60),foreground='#5783a3',background='#1d5075').place(x=520,y=30)

# Text Button7
txtbxvendor = Label(master, text='Vendor                        :').place(x=8,y=190)
txtbxvendor1 = Entry(master, textvariable=VendorN).place(x=150,y=190)

#browse button1
pushbuttnsap = Label(master, text='SAP Package PPT      :').place(x=8,y=250)
# Function to update the Listbox with folder contents
def browse_folder():
    folder_path = filedialog.askdirectory()  # Open a folder selection dialog
    if folder_path:
        folder_contents.delete(0, tk.END)  # Clear the Listbox
        for item in os.listdir(folder_path):
            folder_contents.insert(tk.END, item)  # Insert folder contents into Listbox

# Create a button to browse for a folder
pushbuttnsap1 = tk.Button(master, text="Browse SAP Package", command=browse_folder, activebackground="#1d5075",activeforeground="#8eb9d8",bd=3).place(x=150,y=246)


#browse button2
pushbuttnvendor = Label(master, text='Vendor info PPT        :').place(x=8,y=300)
# Function to update the Listbox with folder contents
def browse_folder():
    folder_path = filedialog.askdirectory()  # Open a folder selection dialog
    if folder_path:
        folder_contents.delete(0, tk.END)  # Clear the Listbox
        for item in os.listdir(folder_path):
            folder_contents.insert(tk.END, item)  # Insert folder contents into Listbox

# Create a button to browse for a folder
pushbuttnvendor1 = tk.Button(master, text="Browse Vendor Info", command=browse_folder).place(x=150,y=296)

#Select List button
selectlist = Label(master, text='Engine Model            :').place(x=8,y=350)
# Create the list of options 
options_list = ["PW1100G", "PW1500G", "PW1900G"] 
  
# Variable to keep track of the option 
# selected in OptionMenu 
value_inside = tk.StringVar(master) 
  
# Set the default value of the variable 
value_inside.set("Select Engine Model") 
  
# Create the optionmenu widget and passing  
# the options_list and value_inside to it. 
question_menu = tk.OptionMenu(master, value_inside, *options_list).place(x=150,y=346)
#question_menu.pack()


# Function to print the submitted option-- testing purpose 
def print_answers(): 
    print("Selected Option: {}".format(value_inside.get())) 
    return None



# Function to update the Listbox with folder contents
def browse_folder():
    folder_path = filedialog.askdirectory()  # Open a folder selection dialog
    if folder_path:
        folder_contents.delete(0, tk.END)  # Clear the Listbox
        for item in os.listdir(folder_path):
            folder_contents.insert(tk.END, item)  # Insert folder contents into Listbox



btnTotal=Button(fg="black",font=('ariel' ,10,'bold'),height=2, width=8, text="OK", command=Ref_ok)
btnTotal.place(x=420,y=496)

btnexit=Button(fg="black",font=('ariel' ,10,'bold'),height=2, width=8, text="Generate", command=qexit)
btnexit.place(x=496,y=496)

# Execute tkinter
master.mainloop()

##################################
####### Summary Table ############
##################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[5]  
mySlide = myPPT.slides.add_slide(slideLayout)  
shapes = mySlide.shapes  

#font = shapes.font
shapes.title.text = 'Summary Table'
#shapes.title.text = 'Prior History by MQI or Damage Code'
title_shape = shapes.title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# Insert shape
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.line.color.rgb = RGBColor(255, 1, 1)
line1.line.width = Pt(3.0)

# specifying the rows, columns, and other measurements  
LI1Val1 = LI1Val.get()
print(LI1Val1)
rows = int(LI1Val1)
rows1 = rows + 1
#rows = 3
cols = 5  
left = Inches(0.5)
top = Inches(2.0)  
width = Inches(9.0)  
height = Inches(0.8)  
  
# using the add_table() method  
myTable = shapes.add_table(rows1, cols, left, top, width, height).table  


  
# writing column headings  
myTable.cell(0, 0).text = 'Line Item'  
myTable.cell(0, 1).text = 'LI'  
myTable.cell(0, 2).text = 'Location'  
myTable.cell(0, 3).text = 'Description'
myTable.cell(0, 4).text = 'Recommendation'  

# writing body cells  
myTable.cell(1, 0).text = '1'  
myTable.cell(1, 1).text = 'Ram'
myTable.cell(1, 2).text = 'Enter manually'
myTable.cell(1, 3).text = 'Enter manually'
myTable.cell(1, 4).text = 'Enter manually'


  
# saving the PPT file  
myPPT.save('Output.pptx')  

#############################################
########## Defect Description ###############
#############################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  
    shapes.title.text = 'Defect Description'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    #rectangle = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
    #shape.text= "Thank You"
    #line=shape.line
    #line.color.rgb=RGBColor(255,0,0)


    right = Inches(8)
    top = Inches(1.10)
    width =  Inches(1.5) 
    height = Inches(0.25)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 


    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"


    #shape = mySlide.Shapes.AddShape(msoShapeRectangle, x, y, w, h)
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2), Inches(2), Inches(2), Inches(2))

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(10) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 
    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item 1: Enter manually"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.9), Inches(0), Inches(1.9))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    # saving the PPT file  
    myPPT.save('Output.pptx')  

    ##############################################
    ########### Defect Location ##################
    ##############################################
    # creating the slide  
    myPPT = Presentation('Output.pptx')  
    #for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Defect Location'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
    #shape.text= "Thank You"
    #line=shape.line
    #line.color.rgb=RGBColor(255,0,0)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0.1)
    top = Inches(1.5)
    width =  Inches(10) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Line Item 1: B/P Location of defect"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.9), Inches(0), Inches(1.9))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)
        

    # saving the PPT file  
    myPPT.save('Output.pptx')  


##########################################
#### Defect MQI Location #############
##########################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Defect MQI Location'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
       
    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
    #shape.text= "Thank You"
    #line=shape.line
    #line.color.rgb=RGBColor(255,0,0)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(10) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Line Item 1 : B/P Location of defect"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.9), Inches(0), Inches(1.9))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    # saving the PPT file  
    myPPT.save('Output.pptx') 


##########################################
#### Disposition & Approvals #############
##########################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Disposition & Approvals'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
#title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = '\n'

# saving the PPT file  
myPPT.save('Output.pptx')  

#####################################################
######### Interim502 Recommendation #################
#####################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font 
    shapes.title.text = 'Interim502 Recommendation XXX'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'

    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
    #shape.text= "Thank You"
    #line=shape.line
    #line.color.rgb=RGBColor(255,0,0)

    right = Inches(4)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "QN XXXXXXXXX "


    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 
    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "LINE ITEM X"



    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(4) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item-1: XXX"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    right = Inches(0)
    top = Inches(1.9)
    width =  Inches(10) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Line Item-1, MQI XXX: "

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)


    # saving the PPT file  
    myPPT.save('Output.pptx')  

#################################################
#########   Design PL3 Approval   ###############
#################################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Design PL3 Approval'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)


    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(4)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "QN XXXXXXXXX "

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(4) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Line Item X"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################
######## Design Assessment  ##############
##########################################
# creating the slide  
myPPT = Presentation('Output.pptx')
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Design Assessment'
  
title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
#title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = '\n'

# saving the PPT file  
myPPT.save('Output.pptx')  


#########################################################
########### Design Analysis Summary ####################
#########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')
for i in range(rows): 
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Design Analysis Summary'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)
    
    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(4) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Line Item - 1: XXX"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    # saving the PPT file  
    myPPT.save('Output.pptx')

########################################################
################  Engine Cross Section  ################
########################################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Engine Cross Section'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item X: Defect location in Engine Cross Section"
    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    right = Inches(4)
    top = Inches(2)
    width =  Inches(2.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "Engine Model: XXXXXX "

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################################
#############  Interfacing Parts  ########################
##########################################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Interfacing Parts'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)


    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item X: Enter Manually"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)

    # saving the PPT file  
    myPPT.save('Output.pptx')

    ###############################################################
    ################ Interfacing Parts (3D View) ##################
    ###############################################################
    # creating the slide  
    myPPT = Presentation('Output.pptx')  
    #for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Interfacing Parts (3D View)'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Line Item-1: Enter Manually"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)
    
    right = Inches(5)
    top = Inches(1.75)
    width =  Inches(4) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Engine Model: XXXXXXXX"
    
    # saving the PPT file  
    myPPT.save('Output.pptx')


######################################################
##########  SPEC: XXXX XXXXXX ########################
######################################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
#for i in range(rows):
slideLayout = myPPT.slide_layouts[5]  
mySlide = myPPT.slides.add_slide(slideLayout)  
shapes = mySlide.shapes  
    
#font = shapes.font
shapes.title.text = 'SPEC: XXX XXXXXX'
#slide.shapes.title.font.name = 'Arial'
#shapes.title.name = 'Arial'
title_shape = shapes.title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

# Insert shape
line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
line1.line.color.rgb = RGBColor(255, 1, 1)
line1.line.width = Pt(3.0)

rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
#rectangle.color.rgb = RGBColor(55, 238, 95)
shape = rectangle
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(14, 127, 38)

right = Inches(8)
top = Inches(1.15)
width =  Inches(1.5) 
height = Inches(0.5)  
    
# creating textBox 
txBox = mySlide.shapes.add_textbox(right, top, width, height) 

# creating textFrames 
tf = txBox.text_frame 
tf.text = "LINE ITEM X"

right = Inches(0)
top = Inches(1.5)
width =  Inches(6) 
height = Inches(0.4)  
    
# creating textBox 
txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

# creating textFrames 
tf1 = txBox1.text_frame 
tf1.text = "Item-1: Enter Manually"

for paragraph in tf1.paragraphs:
    for run in paragraph.runs:
        run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)
    #line2.line.width = Pt(3.0)


right = Inches(0)
top = Inches(1.8)
width =  Inches(6) 
height = Inches(0.4)  

# creating textBox 
txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

# creating textFrames 
tf1 = txBox1.text_frame 
tf1.text = "Item-2: Enter Manually"
for paragraph in tf1.paragraphs:
    for run in paragraph.runs:
        run.font.color.rgb = RGBColor(0, 0, 255)

line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
line2.line.color.rgb = RGBColor(55, 131, 238)
#line2.line.width = Pt(3.0)
    
# saving the PPT file  
myPPT.save('Output.pptx')

###########################################################
#################  Vendor Information  ####################
###########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Vendor Information'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
#title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = '\n'

# saving the PPT file  
myPPT.save('Output.pptx')  

#######################################################
########   Vendor Supplied Information ################
#######################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Vendor Supplied Information'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"
    # saving the PPT file  
    myPPT.save('Output.pptx')

##############################################################
############## Engine Manual   ###############################
##############################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Engine Manual'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
#title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = '\n'

# saving the PPT file  
myPPT.save('Output.pptx')  

########################################################
############ Engine Manual #############################
########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Engine Manual'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item X: Enter Manually"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)

    # saving the PPT file  
    myPPT.save('Output.pptx')

####################################################################
########### Root Cause, Corrective Action & Prior History ##########
####################################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
slideLayout = myPPT.slide_layouts[0]
mySlide = myPPT.slides.add_slide(slideLayout)
title = mySlide.shapes.title
title.text = 'Root Cause, Corrective Action & Prior History'

title_shape = title
title_shape.text_frame.paragraphs[0].font.name = 'Arial'
title_shape.text_frame.paragraphs[0].font.size = Pt(28)
title_shape.text_frame.paragraphs[0].font.bold = True

subtitle = mySlide.placeholders[1]
background = mySlide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(142, 146, 211)
title.bold = True
#title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = '\n'

# saving the PPT file  
myPPT.save('Output.pptx')

#########################################################
########### Root Cause & Corrective Action ##############
#########################################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Root Cause & Corrective Action'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item X: Root cause & Corrective Action "

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)

    # saving the PPT file  
    myPPT.save('Output.pptx')

#########################################################
########### Prior History by S/N  #######################
#########################################################
# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Prior History by S/N'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)


    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item X: Enter Manually"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)

    # saving the PPT file  
    myPPT.save('Output.pptx')

##########################################################
#########  Prior History by MQI or Damage Code ###########
##########################################################

# creating the slide  
myPPT = Presentation('Output.pptx')  
for i in range(rows):
    slideLayout = myPPT.slide_layouts[5]  
    mySlide = myPPT.slides.add_slide(slideLayout)  
    shapes = mySlide.shapes  

    #font = shapes.font
    shapes.title.text = 'Prior History by MQI or Damage Code'
    #slide.shapes.title.font.name = 'Arial'
    #shapes.title.name = 'Arial'
    title_shape = shapes.title
    title_shape.text_frame.paragraphs[0].font.name = 'Arial'
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Insert shape
    line1= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.5), Inches(0), Inches(1.5))
    line1.line.color.rgb = RGBColor(255, 1, 1)
    line1.line.width = Pt(3.0)

    rectangle = mySlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.1),Inches(1.5),Inches(0.35))
    #rectangle.color.rgb = RGBColor(55, 238, 95)
    shape = rectangle
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(14, 127, 38)

    right = Inches(8)
    top = Inches(1.15)
    width =  Inches(1.5) 
    height = Inches(0.5)  
    
    # creating textBox 
    txBox = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf = txBox.text_frame 
    tf.text = "LINE ITEM X"

    right = Inches(0)
    top = Inches(1.5)
    width =  Inches(6) 
    height = Inches(0.4)  
    
    # creating textBox 
    txBox1 = mySlide.shapes.add_textbox(right, top, width, height) 

    # creating textFrames 
    tf1 = txBox1.text_frame 
    tf1.text = "Item X: Enter Manually"

    for paragraph in tf1.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 255)

    line2= mySlide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10), Inches(1.85), Inches(0), Inches(1.85))
    line2.line.color.rgb = RGBColor(55, 131, 238)

    # saving the PPT file  
    myPPT.save('Output.pptx')